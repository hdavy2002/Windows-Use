"""
humphi_search/indexer.py

Fast file indexer for Humphi AI.
Indexes filename, folder, date, size + a tiny content snippet
from every file regardless of size. Never reads more than needed.
"""

import sys
import hashlib
import logging
from pathlib import Path
from datetime import datetime

logging.basicConfig(level=logging.INFO, format="[%(asctime)s] %(message)s", datefmt="%H:%M:%S")
log = logging.getLogger("humphi.indexer")

CHROMA_DIR = Path.home() / ".humphi" / "vectordb"
CHROMA_DIR.mkdir(parents=True, exist_ok=True)

SUPPORTED_DOCS   = {".docx", ".pdf", ".xlsx", ".pptx", ".doc", ".xls"}
SUPPORTED_IMAGES = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".heic", ".tiff"}
SUPPORTED_EXTENSIONS = SUPPORTED_DOCS | SUPPORTED_IMAGES

# Filename patterns to skip
SKIP_NAME_PATTERNS = {
    "~$", ".tmp", "thumbs.db", ".ds_store", "desktop.ini",
    "autorecovery", "autorecover", "~wrd", "~wrl", "corrupt", "cache",
}

# Folder names to skip entirely
SKIP_FOLDERS = {
    "node_modules", ".git", ".vscode", ".idea", "__pycache__",
    "venv", "env", ".env", "dist", "build", "target", "bin", "obj",
    ".next", ".nuxt", "vendor", "bower_components", "packages",
    "site-packages", "AppData", "Temp", "tmp", "$Recycle.Bin",
    "System Volume Information", "Windows", "Program Files",
    "Program Files (x86)", "ProgramData", ".thumbnails", ".cache",
}

SNIPPET_CHARS = 400


def _should_skip(path: Path) -> bool:
    name = path.name.lower()
    if any(p in name for p in SKIP_NAME_PATTERNS):
        return True
    path_parts_lower = {part.lower() for part in path.parts}
    if path_parts_lower & {f.lower() for f in SKIP_FOLDERS}:
        return True
    for part in path.parts:
        if part.startswith(".") and part not in (".", ".."):
            return True
    return False


def _file_hash(path: Path) -> str:
    stat = path.stat()
    return hashlib.md5(f"{stat.st_size}_{stat.st_mtime}".encode()).hexdigest()


def _get_snippet(path: Path) -> str:
    """
    Extract a tiny snippet from any file regardless of size.
    Only reads the absolute minimum needed to understand what the file is about.
    Never reads the whole file.
    """
    ext = path.suffix.lower()

    try:
        # Plain text files — just read first few lines directly
        if ext in {".txt", ".md", ".csv"}:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read(SNIPPET_CHARS)

        # PDF — read first page text only
        elif ext == ".pdf":
            try:
                import pypdf
                reader = pypdf.PdfReader(str(path))
                if reader.pages:
                    text = reader.pages[0].extract_text() or ""
                    return text[:SNIPPET_CHARS]
            except Exception:
                pass

        # Word doc — read first 5 paragraphs only
        elif ext == ".docx":
            try:
                import docx
                doc = docx.Document(str(path))
                text = " ".join(p.text for p in doc.paragraphs[:5] if p.text.strip())
                return text[:SNIPPET_CHARS]
            except Exception:
                pass

        # Excel — read first row of first sheet only (headers)
        elif ext in {".xlsx", ".xls"}:
            try:
                import openpyxl
                wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
                ws = wb.active
                rows = []
                for i, row in enumerate(ws.iter_rows(max_row=3, values_only=True)):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append(", ".join(cells))
                    if i >= 2:
                        break
                wb.close()
                return " | ".join(rows)[:SNIPPET_CHARS]
            except Exception:
                pass

        # PowerPoint — read first slide text only
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                if prs.slides:
                    slide = prs.slides[0]
                    texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                    return " ".join(texts)[:SNIPPET_CHARS]
            except Exception:
                pass

    except Exception:
        pass

    return ""


def _build_document_text(path: Path) -> str:
    """
    Build searchable text = metadata + tiny snippet.
    Metadata is always instant. Snippet is a fixed tiny read.
    """
    ext = path.suffix.lower()
    stem = path.stem.replace("_", " ").replace("-", " ").replace(".", " ")

    # Folder breadcrumb — last 3 folders in path
    parts = path.parts
    folder_context = " > ".join(parts[-4:-1]) if len(parts) > 3 else str(path.parent)

    try:
        stat = path.stat()
        modified = datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d")
        size_kb = stat.st_size // 1024
    except Exception:
        modified = ""
        size_kb = 0

    # Metadata part — always instant
    metadata_text = (
        f"{stem} "
        f"{ext[1:].upper()} "
        f"{folder_context} "
        f"modified {modified} "
        f"{size_kb}KB"
    )

    # Snippet part — tiny fixed read from file
    if ext in SUPPORTED_DOCS:
        snippet = _get_snippet(path)
        if snippet.strip():
            return f"{metadata_text} | {snippet}"

    return metadata_text


class HumphiIndexer:

    def __init__(self):
        self._init_chroma()
        self.indexed = 0
        self.skipped = 0
        self.errors = 0

    def _init_chroma(self):
        try:
            import chromadb
            from chromadb.utils import embedding_functions
            self.client = chromadb.PersistentClient(path=str(CHROMA_DIR))
            self.ef = embedding_functions.SentenceTransformerEmbeddingFunction(
                model_name="all-MiniLM-L6-v2"
            )
            self.collection = self.client.get_or_create_collection(
                name="humphi_files",
                embedding_function=self.ef,
            )
            log.info(f"ChromaDB ready. {self.collection.count()} files already indexed.")
        except ImportError as e:
            log.error(f"Missing: {e}. Run: pip install chromadb sentence-transformers")
            sys.exit(1)

    def _is_indexed(self, path: Path) -> bool:
        try:
            result = self.collection.get(ids=[str(path)], include=["metadatas"])
            if result["ids"]:
                return result["metadatas"][0].get("hash", "") == _file_hash(path)
        except Exception:
            pass
        return False

    def index_file(self, path: Path) -> bool:
        try:
            if _should_skip(path) or not path.exists() or not path.is_file():
                return False
            if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
                return False
            if self._is_indexed(path):
                self.skipped += 1
                return False

            text = _build_document_text(path)
            if not text.strip():
                return False

            stat = path.stat()
            self.collection.upsert(
                ids=[str(path)],
                documents=[text],
                metadatas=[{
                    "path": str(path),
                    "filename": path.name,
                    "extension": path.suffix.lower(),
                    "folder": str(path.parent),
                    "size_bytes": stat.st_size,
                    "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                    "hash": _file_hash(path),
                    "type": "image" if path.suffix.lower() in SUPPORTED_IMAGES else "document",
                    "indexed_at": datetime.now().isoformat(),
                }]
            )
            self.indexed += 1
            return True

        except Exception as e:
            self.errors += 1
            log.debug(f"Error indexing {path.name}: {e}")
            return False

    def scan_folder(self, folder: Path, progress_callback=None):
        folder = Path(folder)
        if not folder.exists():
            log.warning(f"Not found: {folder}")
            return

        files = [
            f for f in folder.rglob("*")
            if f.is_file()
            and f.suffix.lower() in SUPPORTED_EXTENSIONS
            and not _should_skip(f)
        ]

        total = len(files)
        log.info(f"Found {total} files in {folder.name}")

        for i, file_path in enumerate(files):
            self.index_file(file_path)
            if progress_callback:
                progress_callback(i + 1, total, file_path.name)
            elif (i + 1) % 100 == 0 or (i + 1) == total:
                log.info(f"  {i+1}/{total} — indexed={self.indexed} skipped={self.skipped}")

    def remove_file(self, path: Path):
        try:
            self.collection.delete(ids=[str(path)])
        except Exception:
            pass

    def stats(self) -> dict:
        return {
            "total_indexed": self.collection.count(),
            "session_indexed": self.indexed,
            "session_skipped": self.skipped,
            "session_errors": self.errors,
        }